"""Making documents: PDF, Word, Excel, PowerPoint, and the plain formats.

Seven output formats behind **one** tool, and that is the whole design
decision. Selection accuracy falls off past roughly twenty declarations, so
seven separate create_pdf / create_docx / create_xlsx tools would buy
capability by spending the thing that makes the agent pick correctly. One tool
with a `format` enum costs a single slot, and an enum is easier for a model to
get right than a choice between seven similarly-named tools.

Everything writes through the same path guard as the rest of the file tools,
so a document lands in the workspace, Desktop, Documents or Downloads and
nowhere else, and never over a credential file.
"""
from __future__ import annotations

import csv
import io
from pathlib import Path

from .base import ToolError, tool
from .files import _relative, _resolve

FORMATS = ("pdf", "docx", "xlsx", "pptx", "csv", "html", "md", "txt")

# A page of prose is fine; a novel is a sign something has gone wrong.
MAX_CONTENT = 400_000


def _target(path: str, fmt: str) -> Path:
    """Resolve the destination, appending the extension if it is missing."""
    candidate = path if path.lower().endswith(f".{fmt}") else f"{path}.{fmt}"
    target = _resolve(candidate)
    target.parent.mkdir(parents=True, exist_ok=True)
    return target


def _rows_from(content: str) -> list[list[str]]:
    """Parse tabular text.

    Accepts real CSV, and also the pipe tables models like to produce, since
    asking for a spreadsheet and getting back a markdown table is the common
    case rather than the exception.
    """
    lines = [ln for ln in content.splitlines() if ln.strip()]
    if not lines:
        return []

    if lines[0].count("|") >= 2:
        rows = []
        for line in lines:
            # Skip the |---|---| separator row markdown tables carry.
            if set(line.replace("|", "").replace(" ", "")) <= {"-", ":"}:
                continue
            rows.append([c.strip() for c in line.strip().strip("|").split("|")])
        return rows

    return [row for row in csv.reader(io.StringIO(content)) if row]


def _write_pdf(target: Path, content: str, title: str) -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()

    def block(text: str, size: int, bold: bool, height: float) -> None:
        """Write a wrapped paragraph.

        multi_cell(w=0) means "to the right margin", which throws
        "Not enough horizontal space" whenever the cursor is not already at
        the left margin -- and after a previous multi_cell it often is not.
        Resetting x and passing the effective page width explicitly avoids
        depending on where the last call happened to leave the cursor.
        """
        pdf.set_font("Helvetica", "B" if bold else "", size)
        pdf.set_x(pdf.l_margin)
        # The core fonts are Latin-1 only. Losing one character beats
        # refusing to write the document.
        safe = text.encode("latin-1", "replace").decode("latin-1")
        pdf.multi_cell(pdf.epw, height, safe)

    if title:
        block(title, 16, True, 10)
        pdf.ln(2)

    for line in content.splitlines() or [""]:
        stripped = line.strip()
        if stripped.startswith("## "):
            block(stripped[3:], 12, True, 7)
        elif stripped.startswith("# "):
            block(stripped[2:], 14, True, 8)
        elif not stripped:
            pdf.ln(4)
        else:
            block(line, 11, False, 6)

    pdf.output(str(target))


def _write_docx(target: Path, content: str, title: str) -> None:
    from docx import Document

    document = Document()
    if title:
        document.add_heading(title, level=0)
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith("## "):
            document.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            document.add_heading(stripped[2:], level=1)
        elif stripped.startswith(("- ", "* ")):
            document.add_paragraph(stripped[2:], style="List Bullet")
        else:
            document.add_paragraph(line)
    document.save(str(target))


def _write_xlsx(target: Path, content: str, title: str) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    book = Workbook()
    sheet = book.active
    sheet.title = (title or "Sheet1")[:31]  # Excel refuses longer names

    rows = _rows_from(content)
    if not rows:
        raise ToolError(
            "no table found in the content",
            hint="give CSV rows, or a markdown table with | separators",
        )

    for row in rows:
        sheet.append(row)
    for cell in sheet[1]:
        cell.font = Font(bold=True)

    # Width from the content, so the result is readable without dragging.
    for column in sheet.columns:
        longest = max((len(str(c.value or "")) for c in column), default=10)
        sheet.column_dimensions[column[0].column_letter].width = min(longest + 2, 60)
    book.save(str(target))


def _write_pptx(target: Path, content: str, title: str) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    deck = Presentation()
    if title:
        opening = deck.slides.add_slide(deck.slide_layouts[0])
        opening.shapes.title.text = title

    # Blank lines separate slides; each block's first line is its heading.
    for block in [b for b in content.split("\n\n") if b.strip()]:
        lines = [ln for ln in block.splitlines() if ln.strip()]
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = lines[0].lstrip("# ").strip()
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, line in enumerate(lines[1:]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = line.lstrip("-* ").strip()
            paragraph.font.size = Pt(18)
    deck.save(str(target))


def _write_html(target: Path, content: str, title: str) -> None:
    from html import escape

    body = "\n".join(
        f"<p>{escape(line)}</p>" if line.strip() else "<br>"
        for line in content.splitlines()
    )
    target.write_text(
        "<!doctype html>\n<html><head><meta charset='utf-8'>"
        f"<title>{escape(title or target.stem)}</title>"
        "<style>body{font-family:system-ui,sans-serif;max-width:46rem;"
        "margin:3rem auto;padding:0 1rem;line-height:1.6}</style></head><body>"
        + (f"<h1>{escape(title)}</h1>" if title else "")
        + body
        + "</body></html>",
        encoding="utf-8",
    )


@tool(group="documents")
def create_document(
    path: str, content: str, format: str = "pdf", title: str = ""
) -> dict:
    """Create a document: PDF, Word, Excel, PowerPoint, CSV, HTML or text.

    Use this when the user asks for a file they would open in an application
    rather than read in the chat -- a report, a letter, a spreadsheet, a deck.
    For plain notes you will read back later, write_file is simpler.

    Markdown-style structure in the content is understood: "# " and "## " make
    headings, "- " makes bullets. For xlsx give CSV rows or a markdown table.
    For pptx, separate slides with a blank line; each block's first line is
    that slide's title.

    Args:
        path: Where to save it, e.g. "Desktop/report" or "notes/summary.pdf".
            The extension is added if you leave it off.
        content: The text to put in the document.
        format: One of pdf, docx, xlsx, pptx, csv, html, md, txt.
        title: Optional heading, used as the document or sheet title.
    """
    fmt = format.lower().lstrip(".")
    if fmt not in FORMATS:
        raise ToolError(
            f"unknown format: {format}", hint=f"use one of {', '.join(FORMATS)}"
        )
    if not content.strip():
        raise ToolError("nothing to write", hint="give the document some content")
    if len(content) > MAX_CONTENT:
        raise ToolError(
            "that is too much content for one document",
            hint="split it across several files",
        )

    target = _target(path, fmt)

    # Models routinely pass title="Report" and open the content with
    # "# Report", which renders the heading twice. Dropping the duplicate is
    # kinder than telling the model off in the description, and it showed up
    # on the very first real document produced.
    if title.strip():
        lines = content.lstrip().splitlines()
        if lines and lines[0].lstrip("# ").strip().lower() == title.strip().lower():
            content = "\n".join(lines[1:]).lstrip("\n")

    try:
        if fmt == "pdf":
            _write_pdf(target, content, title)
        elif fmt == "docx":
            _write_docx(target, content, title)
        elif fmt == "xlsx":
            _write_xlsx(target, content, title)
        elif fmt == "pptx":
            _write_pptx(target, content, title)
        elif fmt == "html":
            _write_html(target, content, title)
        elif fmt == "csv":
            rows = _rows_from(content)
            with target.open("w", newline="", encoding="utf-8") as handle:
                csv.writer(handle).writerows(rows or [[content]])
        else:  # md, txt
            header = f"# {title}\n\n" if title and fmt == "md" else ""
            target.write_text(header + content, encoding="utf-8")
    except ToolError:
        raise
    except ImportError as exc:
        raise ToolError(
            f"the library for {fmt} is not installed ({exc.name})",
            hint="run: pip install -r requirements.txt",
        ) from None
    except Exception as exc:  # noqa: BLE001
        raise ToolError(
            f"could not write the {fmt}: {type(exc).__name__}: {exc}",
            hint="check the content, or try a simpler format such as txt",
        ) from None

    return {
        "created": _relative(target),
        "format": fmt,
        "size_bytes": target.stat().st_size,
        "full_path": str(target),
    }


@tool(group="documents", untrusted_output=True)
def read_document(path: str, max_chars: int = 8000) -> dict:
    """Read the text out of a PDF, Word document, spreadsheet or deck.

    Use this when the user points at a document read_file cannot open -- it
    only handles plain text. Formatting and images are not preserved; this
    extracts the words.

    Args:
        path: The document, e.g. "Desktop/contract.pdf".
        max_chars: Stop after this much text.
    """
    target = _resolve(path)
    if not target.is_file():
        raise ToolError(
            f"no such file: {_relative(target)}", hint="check with list_directory"
        )

    suffix = target.suffix.lower()
    limit = max(500, min(int(max_chars), 40000))

    try:
        if suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(str(target))
            text = "\n\n".join((p.extract_text() or "") for p in reader.pages[:60])
            extra = {"pages": len(reader.pages)}
        elif suffix == ".docx":
            from docx import Document

            document = Document(str(target))
            text = "\n".join(p.text for p in document.paragraphs)
            extra = {"paragraphs": len(document.paragraphs)}
        elif suffix in (".xlsx", ".xlsm"):
            from openpyxl import load_workbook

            book = load_workbook(str(target), read_only=True, data_only=True)
            chunks = []
            for sheet in book.worksheets:
                chunks.append(f"--- {sheet.title} ---")
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    chunks.append(", ".join("" if c is None else str(c) for c in row))
            text, extra = "\n".join(chunks), {"sheets": len(book.worksheets)}
        elif suffix == ".pptx":
            from pptx import Presentation

            deck = Presentation(str(target))
            chunks = []
            for number, slide in enumerate(deck.slides, 1):
                chunks.append(f"--- slide {number} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        chunks.append(shape.text_frame.text)
            text, extra = "\n".join(chunks), {"slides": len(deck.slides)}
        else:
            raise ToolError(
                f"read_document does not handle {suffix or 'that type'}",
                hint="it reads pdf, docx, xlsx and pptx; use read_file for text",
            )
    except ToolError:
        raise
    except ImportError as exc:
        raise ToolError(
            f"the library for {suffix} is not installed ({exc.name})",
            hint="run: pip install -r requirements.txt",
        ) from None
    except Exception as exc:  # noqa: BLE001
        raise ToolError(
            f"could not read that document: {type(exc).__name__}",
            hint="it may be corrupt, password-protected, or a scan with no text",
        ) from None

    if not text.strip():
        return {
            "path": _relative(target),
            "text": "",
            "note": "no text found -- this may be a scan rather than real text",
            **extra,
        }

    return {
        "path": _relative(target),
        "text": text[:limit],
        "truncated": len(text) > limit,
        **extra,
    }
